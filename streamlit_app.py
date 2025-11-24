import streamlit as st
import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 页面配置
st.set_page_config(
    page_title="九宫格潜力展示系统",
    page_icon="📊",
    layout="wide"
)

# 列名映射配置（可根据实际Excel文件调整）
COLUMN_MAPPING = {
    # Excel列名 -> 系统内部列名
    '日常盘点-九宫格结果（需输入）': '档位',
    '员工姓名': '姓名',
    '当前组织': '当前组织',
    '日常盘点-梯队结果（需输入）': '梯队',
    '管理职级': '管理职级',
    '专业职级': '专业职级'
}

# 部门-负责人映射表（用于真实数据上传时自动匹配）
DEPARTMENT_SUPERVISOR_MAPPING = {
    # leon(郭凯天)
    '职能线办公室': 'leon(郭凯天)',
    
    # samxu(徐炎)
    '反洗钱与制裁合规部': 'samxu(徐炎)',
    '金融法律合规部': 'samxu(徐炎)',
    '法务综合部': 'samxu(徐炎)',
    '知识产权部': 'samxu(徐炎)',
    '数据合规与隐私保护部': 'samxu(徐炎)',
    '海外法务中心': 'samxu(徐炎)',
    '合规策略中心': 'samxu(徐炎)',
    'IEG法务部': 'samxu(徐炎)',
    '法务平台部': 'samxu(徐炎)',
    'CDG法务部': 'samxu(徐炎)',
    'CSIG法务部': 'samxu(徐炎)',
    'PCG法务部': 'samxu(徐炎)',
    
    # dega(刘夏耘)
    '行政部': 'dega(刘夏耘)',
    '基建部': 'dega(刘夏耘)',
    '集团采购管理部': 'dega(刘夏耘)',
    'IEG公共事务部': 'dega(刘夏耘)',
    '腾讯华东总部': 'dega(刘夏耘)',
    
    # snailcai(蔡光忠)
    '大湾区公共事务部': 'snailcai(蔡光忠)',
    '腾讯西南总部': 'snailcai(蔡光忠)',
    
    # leolyliu(刘勇)
    '腾讯北京总部': 'leolyliu(刘勇)',
    '集团公共事务部': 'leolyliu(刘勇)',
    'CSIG公共事务部': 'leolyliu(刘勇)',
    
    # cyberchen(陈勇)
    'PCG公共事务部': 'cyberchen(陈勇)',
    'WXG公共事务部': 'cyberchen(陈勇)',
    '文化与内容公共事务部': 'cyberchen(陈勇)',
    '数字舆情部': 'cyberchen(陈勇)',
    
    # lucazhu(朱劲松)
    '安全管理部': 'lucazhu(朱劲松)',
    '信息安全发展部': 'lucazhu(朱劲松)',
    
    # jasonsi(司晓)
    '反垄断合规部': 'jasonsi(司晓)',
    '公共战略研究部': 'jasonsi(司晓)',
    '市场经营合规部': 'jasonsi(司晓)'
}

# 自定义CSS样式
st.markdown("""
<style>
    .main-title {
        font-size: 24px;
        font-weight: bold;
        color: #1f4e79;
        margin-bottom: 20px;
    }
    
    .stats-table {
        border: 1px solid #4472c4;
        border-radius: 5px;
        padding: 10px;
        background-color: #e7f3ff;
    }
    
    .grid-cell {
        border: 2px solid #4472c4;
        border-radius: 8px;
        padding: 15px;
        margin: 8px 4px;
        min-height: 220px;
        max-height: 280px;
        background-color: #f8f9fa;
        display: flex;
        flex-direction: column;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        position: relative;
    }
    
    .cell-title {
        font-weight: bold;
        color: #1f4e79;
        font-size: 14px;
        margin-bottom: 8px;
        border-bottom: 1px solid #d0d0d0;
        padding-bottom: 5px;
        background-color: #f8f9fa;
        position: sticky;
        top: 0;
        z-index: 10;
        flex-shrink: 0;
    }
    
    .cell-content {
        font-size: 12px;
        line-height: 1.4;
        flex-grow: 1;
        overflow-y: auto;
        color: #333;
        max-height: 200px;
        padding-right: 5px;
    }
    
    .cell-content::-webkit-scrollbar {
        width: 4px;
    }
    
    .cell-content::-webkit-scrollbar-track {
        background: #f1f1f1;
        border-radius: 2px;
    }
    
    .cell-content::-webkit-scrollbar-thumb {
        background: #c1c1c1;
        border-radius: 2px;
    }
    
    .cell-content::-webkit-scrollbar-thumb:hover {
        background: #a1a1a1;
    }
    
    .grid-row {
        margin-bottom: 10px;
    }
    
    .grid-container {
        position: relative;
        display: flex;
        flex-direction: column;
        align-items: center;
        margin: 20px 0;
    }
    
    .grid-with-axes {
        position: relative;
        display: inline-block;
        margin-left: 60px;
        margin-bottom: 60px;
    }
    
    .y-axis {
        position: absolute;
        left: -50px;
        top: 0;
        height: 100%;
        display: flex;
        flex-direction: column;
        justify-content: center;
        align-items: center;
        writing-mode: vertical-lr;
        text-orientation: mixed;
    }
    
    .y-axis-label {
        font-weight: bold;
        color: #1f4e79;
        font-size: 16px;
        margin-bottom: 10px;
        writing-mode: vertical-lr;
        text-orientation: mixed;
    }
    
    .y-axis-arrow {
        width: 0;
        height: 0;
        border-left: 8px solid transparent;
        border-right: 8px solid transparent;
        border-bottom: 15px solid #1f4e79;
        margin-bottom: 5px;
    }
    
    .x-axis {
        position: absolute;
        bottom: -50px;
        left: 0;
        width: 100%;
        display: flex;
        justify-content: center;
        align-items: center;
    }
    
    .x-axis-label {
        font-weight: bold;
        color: #1f4e79;
        font-size: 16px;
        margin-left: 10px;
    }
    
    .x-axis-arrow {
        width: 0;
        height: 0;
        border-top: 8px solid transparent;
        border-bottom: 8px solid transparent;
        border-left: 15px solid #1f4e79;
        margin-left: 5px;
    }
    
</style>
""", unsafe_allow_html=True)

def load_and_validate_data(uploaded_file, sheet_name="Sheet1"):
    """
    加载并验证Excel数据
    """
    try:
        # 读取Excel文件
        df = pd.read_excel(uploaded_file, sheet_name=sheet_name)
        
        # 获取必需的Excel列名（从配置中获取）
        required_excel_columns = [col for col in COLUMN_MAPPING.keys() 
                                 if col in ['九宫格', '员工姓名', '当前组织', '梯队']]
        missing_columns = [col for col in required_excel_columns if col not in df.columns]
        
        if missing_columns:
            st.error(f"缺少必需的列: {missing_columns}")
            available_columns = list(df.columns)
            st.info(f"Excel文件中的可用列: {', '.join(available_columns)}")
            st.info(f"请确保Excel文件包含以下列: {', '.join(required_excel_columns)}")
            return None
        
        # 使用配置化的列名映射进行重命名
        rename_mapping = {}
        for excel_col, internal_col in COLUMN_MAPPING.items():
            if excel_col in df.columns:
                rename_mapping[excel_col] = internal_col
        
        df = df.rename(columns=rename_mapping)
        st.info(f"已映射列名: {', '.join([f'{k}→{v}' for k, v in rename_mapping.items()])}")
        
        # 从"当前组织"列中提取部门名称
        def extract_department(org_string):
            """从当前组织字符串中提取部门名称"""
            if pd.isna(org_string) or org_string == "":
                return "未知部门"
            
            try:
                # 分割字符串，取第一个"/"到第二个"/"之间的内容
                parts = str(org_string).split('/')
                if len(parts) >= 2:
                    return parts[1].strip()  # 第二部分就是部门名
                else:
                    return "未知部门"
            except:
                return "未知部门"
        
        # 提取部门名称
        df['部门'] = df['当前组织'].apply(extract_department)
        st.info(f"已从'当前组织'列中提取部门信息，共提取{len(df[df['部门'] != '未知部门'])}条有效部门记录")
        
        # 自动映射负责人（基于部门-负责人映射表）
        if '负责人' not in df.columns:
            df['负责人'] = df['部门'].map(DEPARTMENT_SUPERVISOR_MAPPING).fillna('待分配')
            matched_count = len(df[df['负责人'] != '待分配'])
            st.success(f"已根据部门自动匹配负责人，共匹配{matched_count}条记录")
            
            # 显示未匹配的部门
            unmatched_depts = df[df['负责人'] == '待分配']['部门'].unique()
            if len(unmatched_depts) > 0:
                st.warning(f"以下部门未在映射表中找到对应负责人：{', '.join(unmatched_depts)}")
        
        # 不移除任何数据，保持原始数据完整性
        return df
        
    except Exception as e:
        st.error(f"读取Excel文件时出错: {str(e)}")
        return None

def get_potential_level(rating):
    """
    根据档位获取潜力级别
    """
    if rating in [7, 8, 9]:
        return "高潜力"
    elif rating in [4, 5, 6]:
        return "中潜力"
    elif rating in [1, 2, 3]:
        return "低潜力"
    else:
        return "未知"

def extract_professional_level(text):
    """
    从专业职级文本中提取数字等级
    例如："行政建筑类－AA/行政/S10" -> 10
    """
    import re
    if pd.isna(text) or text == "":
        return 0
    # 提取S后面的数字，如S10, S12, S15等
    match = re.search(r'S(\d+)', str(text))
    return int(match.group(1)) if match else 0

def is_valid_tier(tier_value):
    """
    判断是否为有效梯队
    """
    valid_tiers = ["前5%", "前5%-15%", "前15%", "前15%-40%", "前40%", "末5%"]
    return tier_value in valid_tiers

def extract_professional_display(text):
    """
    从专业职级文本中提取显示格式
    例如：'行政建筑类－AA/行政/S10' -> 'S10'
    """
    import re
    if pd.isna(text) or text == "":
        return ""
    # 提取字母+数字的模式，如S10, AA12等
    match = re.search(r'([A-Z])(\d+)', str(text))
    return f"{match.group(1)}{match.group(2)}" if match else ""

def get_tier_focus_data(df):
    """
    获取梯队重点关注的数据分析
    """
    if df is None or df.empty:
        return pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), ""
    
    # 确保必要的列存在，如果不存在则创建空列
    if '管理职级' not in df.columns:
        df['管理职级'] = ""
    if '专业职级' not in df.columns:
        df['专业职级'] = ""
    
    # 1. 基干未盘入梯队：有管理职级但梯队无效
    cadre_not_tiered = df[
        (df['管理职级'].notna()) & 
        (df['管理职级'] != "") & 
        (~df['梯队'].apply(is_valid_tier))
    ].copy()
    
    # 2. 专家末5%：管理职级为空且专业职级≥12且梯队为"末5%"
    expert_bottom5 = df[
        ((df['管理职级'].isna()) | (df['管理职级'] == "")) &
        (df['专业职级'].apply(extract_professional_level) >= 12) &
        (df['梯队'] == "末5%")
    ].copy()
    
    # 3. 专家未盘入梯队：管理职级为空且专业职级≥12但梯队无效
    expert_not_tiered = df[
        ((df['管理职级'].isna()) | (df['管理职级'] == "")) &
        (df['专业职级'].apply(extract_professional_level) >= 12) &
        (~df['梯队'].apply(is_valid_tier))
    ].copy()
    
    # 计算统计信息
    total_cadres = len(df[(df['管理职级'].notna()) & (df['管理职级'] != "")])
    cadre_not_tiered_count = len(cadre_not_tiered)
    expert_bottom5_count = len(expert_bottom5)
    expert_not_tiered_count = len(expert_not_tiered)
    
    cadre_ratio = (cadre_not_tiered_count / total_cadres * 100) if total_cadres > 0 else 0
    
    # 生成标题文字
    title_text = f"共有{cadre_not_tiered_count}名基干未盘入梯队，占比{cadre_ratio:.1f}%；非管理干部专家中盘入末5%共{expert_bottom5_count}人，12级以上专家未盘入梯队共{expert_not_tiered_count}人。"
    
    return cadre_not_tiered, expert_bottom5, expert_not_tiered, title_text

def extract_professional_level(text):
    """
    从专业职级文本中提取数字等级
    例如："行政建筑类－AA/行政/S10" -> 10
    """
    import re
    if pd.isna(text) or text == "":
        return 0
    # 提取S后面的数字，如S10, S12, S15等
    match = re.search(r'S(\d+)', str(text))
    return int(match.group(1)) if match else 0

def is_valid_tier(tier_value):
    """
    判断是否为有效梯队
    """
    valid_tiers = ["前5%", "前5%-15%", "前15%", "前15%-40%", "前40%", "末5%"]
    return tier_value in valid_tiers

def get_tier_focus_data(df):
    """
    获取梯队重点关注的数据分析
    """
    if df is None or df.empty:
        return {}, {}, {}, ""
    
    # 确保必要的列存在，如果不存在则创建空列
    if '管理职级' not in df.columns:
        df['管理职级'] = ""
    if '专业职级' not in df.columns:
        df['专业职级'] = ""
    
    # 1. 基干未盘入梯队：有管理职级但梯队无效
    cadre_not_tiered = df[
        (df['管理职级'].notna()) & 
        (df['管理职级'] != "") & 
        (~df['梯队'].apply(is_valid_tier))
    ].copy()
    
    # 2. 专家末5%：管理职级为空且专业职级≥12且梯队为"末5%"
    expert_bottom5 = df[
        ((df['管理职级'].isna()) | (df['管理职级'] == "")) &
        (df['专业职级'].apply(extract_professional_level) >= 12) &
        (df['梯队'] == "末5%")
    ].copy()
    
    # 3. 专家未盘入梯队：管理职级为空且专业职级≥12但梯队无效
    expert_not_tiered = df[
        ((df['管理职级'].isna()) | (df['管理职级'] == "")) &
        (df['专业职级'].apply(extract_professional_level) >= 12) &
        (~df['梯队'].apply(is_valid_tier))
    ].copy()
    
    # 计算统计信息
    total_cadres = len(df[(df['管理职级'].notna()) & (df['管理职级'] != "")])
    cadre_not_tiered_count = len(cadre_not_tiered)
    expert_bottom5_count = len(expert_bottom5)
    expert_not_tiered_count = len(expert_not_tiered)
    
    cadre_ratio = (cadre_not_tiered_count / total_cadres * 100) if total_cadres > 0 else 0
    
    # 生成标题文字
    title_text = f"共有{cadre_not_tiered_count}名基干未盘入梯队，占比{cadre_ratio:.1f}%；非管理干部专家中盘入末5%共{expert_bottom5_count}人，12级以上专家未盘入梯队共{expert_not_tiered_count}人。"
    
    return cadre_not_tiered, expert_bottom5, expert_not_tiered, title_text

def generate_summary(df, rating):
    """
    生成指定档位的员工汇总信息
    """
    if df is None or df.empty:
        return "无数据", 0
    
    # 筛选指定档位的数据
    rating_data = df[df['档位'] == rating]
    
    if rating_data.empty:
        return "无符合员工", 0
    
    # 按部门分组
    dept_groups = rating_data.groupby('部门')['姓名'].apply(list).to_dict()
    
    # 格式化输出
    summary_lines = []
    total_count = 0
    
    for dept, names in dept_groups.items():
        names_str = "、".join(names)
        summary_lines.append(f"{dept}：{names_str}")
        total_count += len(names)
    
    summary = "\n".join(summary_lines)
    return summary, total_count

def create_stats_table(df):
    """
    创建九宫格统计汇总表格
    """
    if df is None or df.empty:
        return pd.DataFrame()
    
    # 计算九宫格统计
    high_potential_count = len(df[df['档位'].isin([6, 8, 9])])  # 6/8/9
    low_potential_count = len(df[df['档位'].isin([1, 2, 4])])   # 1/2/4
    
    # 实际指标人数
    high_target = 42  # 6/8/9指标
    low_target = 14   # 1/2/4指标
    
    # 计算超额/少打情况
    high_diff = high_potential_count - high_target
    low_diff = low_potential_count - low_target
    
    high_remark = f"超额{high_diff}人" if high_diff > 0 else f"少打{abs(high_diff)}人" if high_diff < 0 else "达标"
    low_remark = f"超额{low_diff}人" if low_diff > 0 else f"少打{abs(low_diff)}人" if low_diff < 0 else "达标"

    # 构建统计表格
    stats_df = pd.DataFrame({
        '九宫格': ['6/8/9', '1/2/4'],
        '指标': [high_target, low_target],
        '已打': [high_potential_count, low_potential_count],
        '备注': [high_remark, low_remark]
    })
    
    return stats_df

def create_tier_stats_table(df):
    """
    创建梯队统计汇总表格
    """
    if df is None or df.empty:
        return pd.DataFrame()
    
    # 梯队映射 - 保留末5%档位
    tier_mapping = {
        # 前5%档：只包含"前5%"
        '前5%': ['前5%'],
        # 前15%档：包含"前5%-15%"和"前15%"
        '前15%': ['前5%-15%', '前15%'],
        # 前40%档：包含"前15%-40%"和"前40%"  
        '前40%': ['前15%-40%', '前40%'],
        # 末5%档：只包含"末5%"
        '末5%': ['末5%']
    }
    
    # 实际指标人数
    tier_targets = {
        '前5%': 72,
        '前15%': 139,
        '前40%': 351,
        '末5%': 72
    }
    
    # 计算各档实际人数
    stats_data = []
    
    # 调试：检查所有唯一的梯队值
    unique_tiers = df['梯队'].unique()
    
    for tier_group, tier_values in tier_mapping.items():
        # 统计该档下所有梯队的人数
        actual_count = 0
        for tier_value in tier_values:
            count = len(df[df['梯队'] == tier_value])
            actual_count += count
        
        target_count = tier_targets.get(tier_group, 0)
        
        stats_data.append({
            '梯队': tier_group,
            '指标': int(target_count),
            '已打': int(actual_count)
        })
    
    stats_df = pd.DataFrame(stats_data)
    # 确保数据类型正确
    stats_df['指标'] = stats_df['指标'].astype(int)
    stats_df['已打'] = stats_df['已打'].astype(int)
    return stats_df

def generate_tier_summary(df):
    """
    生成梯队总览汇总文字
    """
    if df is None or df.empty:
        return "无数据"
    
    tier_stats_df = create_tier_stats_table(df)
    summary_lines = []
    
    # 获取各梯队数据
    tier_data = {}
    for _, row in tier_stats_df.iterrows():
        tier_data[row['梯队']] = {
            'target': row['指标'],
            'actual': row['已打'],
            'diff': row['已打'] - row['指标']
        }
    
    # 计算前三档总数
    top_three_target = tier_data.get('前5%', {}).get('target', 0) + \
                      tier_data.get('前15%', {}).get('target', 0) + \
                      tier_data.get('前40%', {}).get('target', 0)
    
    top_three_actual = tier_data.get('前5%', {}).get('actual', 0) + \
                      tier_data.get('前15%', {}).get('actual', 0) + \
                      tier_data.get('前40%', {}).get('actual', 0)
    
    # 检查总体约束
    total_violation = top_three_actual > top_three_target

    # 按梯队顺序生成汇总
    tier_order = ['前5%', '前15%', '前40%', '末5%']
    
    for tier in tier_order:
        if tier not in tier_data:
            continue
            
        data = tier_data[tier]
        target = data['target']
        actual = data['actual']
        diff = data['diff']
        
        # 根据梯队核算逻辑判断状态
        if tier == '前5%':
            # 前5%：可以少于指标，但不能多于指标
            if diff > 0:
                status = f"超额{diff}人"  # 违规
            else:
                status = "达标"
        elif tier == '前15%':
            # 前15%：可以少于指标，但不能多于指标
            if diff > 0:
                status = f"超额{diff}人"  # 违规
            else:
                status = "达标"
        elif tier == '前40%':
            # 前40%：可以多于指标，但需要检查总体约束
            if total_violation:
                status = f"总体超额{top_three_actual - top_three_target}人"
            else:
                status = "达标"
        elif tier == '末5%':
            # 末5%：可以多于指标，但不能少于指标
            if diff < 0:
                status = f"少打{abs(diff)}人"  # 违规
            else:
                status = "达标"
        else:
            status = "达标"
            
        summary_lines.append(f"{tier}{status}")

    return "，".join(summary_lines)

def generate_grid_summary(df):
    """
    生成基干九宫格汇总文字
    """
    if df is None or df.empty:
        return "无数据"
    
    # 只统计有九宫格数据的人数
    total_count = len(df[df['档位'].notna()])
    
    # 6/8/9统计
    high_count = len(df[df['档位'].isin([6, 8, 9])])
    high_target = 42
    high_diff = high_count - high_target
    high_status = f"超额{high_diff}人" if high_diff > 0 else f"少打{abs(high_diff)}人" if high_diff < 0 else "达标"
    
    # 1/2/4统计
    low_count = len(df[df['档位'].isin([1, 2, 4])])
    low_target = 14
    low_diff = low_count - low_target
    low_status = f"超额{low_diff}人" if low_diff > 0 else f"少打{abs(low_diff)}人" if low_diff < 0 else "达标"
    
    summary = f"总计{total_count}人，6/8/9{high_status}，1/2/4{low_status}"
    
    return summary

def create_sample_data():
    """
    创建示例数据用于测试 - 1500人，32个部门
    """
    import random
    import numpy as np
    
    # 设置随机种子确保结果可重现
    random.seed(42)
    np.random.seed(42)
    
    # 32个部门
    departments = [
        '技术研发部', '产品设计部', '市场营销部', '销售部', '人力资源部', '财务部', '法务部', '行政部',
        '运营部', '客服部', '品牌部', '商务拓展部', '数据分析部', '质量管理部', '项目管理部', '战略规划部',
        '供应链部', '采购部', '生产部', '研发工程部', '测试部', '安全部', '合规部', '投资部',
        '公关部', '培训部', 'IT支持部', '用户体验部', '内容运营部', '商业智能部', '风控部', '审计部'
    ]
    
    # 8个负责人
    supervisors = ['张总监', '李经理', '王主管', '陈总', '刘部长', '赵经理', '孙主任', '周总监']
    
    # 常见姓氏和名字
    surnames = ['张', '李', '王', '刘', '陈', '杨', '赵', '黄', '周', '吴', '徐', '孙', '胡', '朱', '高', '林', '何', '郭', '马', '罗', '梁', '宋', '郑', '谢', '韩', '唐', '冯', '于', '董', '萧', '程', '曹', '袁', '邓', '许', '傅', '沈', '曾', '彭', '吕']
    given_names = ['伟', '芳', '娜', '秀英', '敏', '静', '丽', '强', '磊', '军', '洋', '勇', '艳', '杰', '娟', '涛', '明', '超', '秀兰', '霞', '平', '刚', '桂英', '建华', '文', '华', '志强', '秀珍', '春梅', '海燕', '雪', '建国', '建军', '晓东', '梅', '丹', '雨', '辉', '玲', '燕']
    
    # 生成1500个员工
    employees = []
    
    # 管理职级分配：120人（L1: 80人，L2: 40人）
    management_levels = ['L1'] * 80 + ['L2'] * 40
    
    # 专业职级分配：1380人（S5-S15，其中S12+: 150人）
    professional_levels = []
    # S12-S15: 150人
    professional_levels.extend(['S12'] * 60 + ['S13'] * 40 + ['S14'] * 30 + ['S15'] * 20)
    # S5-S11: 1230人
    professional_levels.extend(['S5'] * 200 + ['S6'] * 200 + ['S7'] * 180 + ['S8'] * 180 + 
                              ['S9'] * 160 + ['S10'] * 160 + ['S11'] * 150)
    
    # 九宫格分布：只在120个管理职级人员中分布
    # 6/8/9=20%(24人), 1/2/4=10%(12人), 3/5/7=70%(84人)
    grid_distribution = []
    # 高潜力 6/8/9: 24人
    grid_distribution.extend([6] * 8 + [8] * 8 + [9] * 8)
    # 低潜力 1/2/4: 12人
    grid_distribution.extend([1] * 4 + [2] * 4 + [4] * 4)
    # 中等 3/5/7: 84人
    grid_distribution.extend([3] * 28 + [5] * 28 + [7] * 28)
    
    # 为专业职级人员准备空的九宫格值
    empty_grid = [None] * 1380  # 专业职级人员没有九宫格评级
    
    # 梯队分布 - 修正版本
    tier_distribution = []
    tier_distribution.extend(['前5%'] * 75)      # 5% - 顶尖人才
    tier_distribution.extend(['前5%-15%'] * 150)  # 10% - 核心人才(5%-15%)
    tier_distribution.extend(['前15%-40%'] * 375) # 25% - 骨干人才(15%-40%)
    tier_distribution.extend(['末5%'] * 75)      # 5% - 重点关注
    tier_distribution.extend([''] * 825)         # 55% - 普通员工(无特殊梯队标识)
    
    # 打乱分布
    random.shuffle(management_levels)
    random.shuffle(professional_levels)
    random.shuffle(grid_distribution)
    random.shuffle(tier_distribution)
    
    # 生成员工数据
    for i in range(1500):
        # 生成姓名
        surname = random.choice(surnames)
        given_name = random.choice(given_names)
        if random.random() < 0.3:  # 30%概率生成两字名
            given_name += random.choice(given_names)
        name = surname + given_name
        
        # 分配部门（确保每个部门都有人）
        if i < 32:
            dept = departments[i]
        else:
            dept = random.choice(departments)
        
        # 分配负责人
        supervisor = random.choice(supervisors)
        
        # 分配职级
        if i < 120:  # 前120人为管理职级
            mgmt_level = management_levels[i]
            prof_level = ''
            prof_level_display = ''
        else:  # 其余为专业职级
            mgmt_level = ''
            prof_level_num = professional_levels[i-120]
            # 根据部门生成专业职级显示格式
            if '技术' in dept or '研发' in dept or 'IT' in dept or '测试' in dept:
                prof_level = f'技术类-AA/技术/{prof_level_num}'
            elif '市场' in dept or '销售' in dept or '商务' in dept:
                prof_level = f'市场类-BB/市场/{prof_level_num}'
            elif '财务' in dept or '审计' in dept:
                prof_level = f'财务类-CC/财务/{prof_level_num}'
            elif '人力' in dept or '培训' in dept:
                prof_level = f'人事类-DD/人事/{prof_level_num}'
            elif '运营' in dept or '客服' in dept:
                prof_level = f'运营类-EE/运营/{prof_level_num}'
            else:
                prof_level = f'综合类-FF/综合/{prof_level_num}'
        
        # 分配九宫格和梯队
        if i < 120:  # 管理职级人员才有九宫格评级
            grid_pos = grid_distribution[i]
        else:  # 专业职级人员没有九宫格评级
            grid_pos = None
        tier = tier_distribution[i]
        
        # 生成备注（20%概率有备注）
        remarks = ['', '重点关注', '优秀员工', '需培训', '待观察', '核心骨干', '新员工', '资深员工']
        remark = random.choice(remarks) if random.random() < 0.2 else ''
        
        employees.append({
            '九宫格': grid_pos,
            '员工姓名': name,
            '部门': dept,
            '梯队': tier,
            '管理职级': mgmt_level,
            '专业职级': prof_level,
            '负责人': supervisor,
            '备注': remark
        })
    
    df = pd.DataFrame(employees)
    # 重命名列以保持代码兼容性
    df = df.rename(columns={'九宫格': '档位', '员工姓名': '姓名'})
    return df

def create_ppt_report(df):
    """
    创建腾讯风格PPT报告
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
    
    # 创建新的演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 腾讯配色方案
    tencent_blue_standard = RGBColor(33, 81, 209)  # 腾讯标准蓝 #2151D1
    tencent_blue_medium = RGBColor(70, 130, 180)  # 中蓝色 #4682B4
    tencent_blue_light = RGBColor(135, 206, 235)  # 浅蓝色 #87CEEB
    tencent_yellow = RGBColor(255, 215, 0)  # 黄色高亮 #FFD700
    white_color = RGBColor(255, 255, 255)
    
    # 第一页：腾讯风格封面页
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景为腾讯标准蓝
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = tencent_blue_standard
    
    # Tencent 腾讯 logo文字
    logo_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6), Inches(1))
    logo_frame = logo_box.text_frame
    logo_p = logo_frame.paragraphs[0]
    logo_p.text = "Tencent 腾讯"
    logo_p.font.size = Pt(28)
    logo_p.font.color.rgb = white_color
    logo_p.font.name = "Microsoft YaHei"
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    current_year = pd.Timestamp.now().year
    current_month = pd.Timestamp.now().strftime('%m')
    title_p.text = f"{current_year}H{1 if int(current_month) <= 6 else 2}"
    title_p.font.size = Pt(72)
    title_p.font.bold = True
    title_p.font.color.rgb = white_color
    title_p.font.name = "Microsoft YaHei"
    
    # 副标题
    subtitle_p = title_frame.add_paragraph()
    subtitle_p.text = f"S1人才盘点情况"
    subtitle_p.font.size = Pt(48)
    subtitle_p.font.color.rgb = white_color
    subtitle_p.font.name = "Microsoft YaHei"
    
    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(4), Inches(0.8))
    date_frame = date_box.text_frame
    date_p = date_frame.paragraphs[0]
    date_p.text = pd.Timestamp.now().strftime('%Y.%m')
    date_p.font.size = Pt(24)
    date_p.font.color.rgb = white_color
    date_p.font.name = "Microsoft YaHei"
    
    # 第二页：数据总览页
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "1 数据总览"
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = tencent_blue_standard
    title_p.font.name = "Microsoft YaHei"
    
    # 创建统计表格
    stats_df = create_stats_table(df)
    tier_stats_df = create_tier_stats_table(df)
    
    # 生成汇总文字
    tier_summary = generate_tier_summary(df)
    grid_summary = generate_grid_summary(df)
    
    # 左侧：梯队总览标题和汇总文字
    tier_title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(4.5), Inches(0.3))
    tier_title_frame = tier_title_box.text_frame
    tier_title_p = tier_title_frame.paragraphs[0]
    tier_title_p.text = "梯队总览"
    tier_title_p.font.size = Pt(16)
    tier_title_p.font.bold = True
    tier_title_p.font.color.rgb = tencent_blue_standard
    tier_title_p.font.name = "Microsoft YaHei"
    
    # 梯队汇总文字
    tier_summary_box = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(4.5), Inches(0.2))
    tier_summary_frame = tier_summary_box.text_frame
    tier_summary_p = tier_summary_frame.paragraphs[0]
    tier_summary_p.text = tier_summary
    tier_summary_p.font.size = Pt(10)
    tier_summary_p.font.color.rgb = RGBColor(102, 102, 102)
    tier_summary_p.font.name = "Microsoft YaHei"
    
    # 左侧：梯队总览表格
    table_rows = len(tier_stats_df) + 1  # 数据行数 + 表头
    table_left = slide.shapes.add_table(table_rows, 3, Inches(1), Inches(1.7), Inches(4.5), Inches(2.8))
    table_left.table.cell(0, 0).text = "梯队"
    table_left.table.cell(0, 1).text = "指标"
    table_left.table.cell(0, 2).text = "已打"
    
    # 设置表头样式
    for i in range(3):
        cell = table_left.table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = tencent_blue_medium
        cell.text_frame.paragraphs[0].font.color.rgb = white_color
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    
    # 填充梯队数据
    for i, (_, row) in enumerate(tier_stats_df.iterrows(), 1):
        table_left.table.cell(i, 0).text = row['梯队']
        table_left.table.cell(i, 1).text = str(row['指标'])
        table_left.table.cell(i, 2).text = str(row['已打'])
        
        for j in range(3):
            cell = table_left.table.cell(i, j)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    
    # 右侧：基干九宫格标题和汇总文字
    grid_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.2), Inches(5.5), Inches(0.3))
    grid_title_frame = grid_title_box.text_frame
    grid_title_p = grid_title_frame.paragraphs[0]
    grid_title_p.text = "基干九宫格"
    grid_title_p.font.size = Pt(16)
    grid_title_p.font.bold = True
    grid_title_p.font.color.rgb = tencent_blue_standard
    grid_title_p.font.name = "Microsoft YaHei"
    
    # 九宫格汇总文字
    grid_summary_box = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(0.2))
    grid_summary_frame = grid_summary_box.text_frame
    grid_summary_p = grid_summary_frame.paragraphs[0]
    grid_summary_p.text = grid_summary
    grid_summary_p.font.size = Pt(10)
    grid_summary_p.font.color.rgb = RGBColor(102, 102, 102)
    grid_summary_p.font.name = "Microsoft YaHei"
    
    # 右侧：基干九宫格统计表
    table_right = slide.shapes.add_table(3, 4, Inches(7), Inches(1.7), Inches(5.5), Inches(2.3))
    
    # 表头
    headers = ["九宫格", "指标", "已打", "备注"]
    for i, header in enumerate(headers):
        cell = table_right.table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = tencent_blue_medium
        cell.text_frame.paragraphs[0].font.color.rgb = white_color
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    
    # 使用统一的统计函数计算九宫格数据
    stats_df = create_stats_table(df)
    
    # 动态生成表格数据
    grid_data = []
    for _, row in stats_df.iterrows():
        grid_type = row['九宫格']
        target = row['指标']
        actual = row['已打']
        remark = row['备注']
        
        grid_data.append([
            grid_type,
            str(target),
            str(actual),
            remark
        ])
    
    for i, row_data in enumerate(grid_data, 1):
        for j, cell_data in enumerate(row_data):
            cell = table_right.table.cell(i, j)
            cell.text = cell_data
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
            
            # 黄色高亮"已打"列
            if j == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = tencent_yellow
    

    
    # 第三页：九宫格详细展示
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    # 计算有九宫格数据的人数
    grid_count = len(df[df['档位'].notna()])
    title_p.text = f"2 基干九宫格一览：总计{grid_count}人"
    title_p.font.size = Pt(20)
    title_p.font.bold = True
    title_p.font.color.rgb = tencent_blue_standard
    title_p.font.name = "Microsoft YaHei"
    
    # 右上角统计表格
    stats_table = slide.shapes.add_table(3, 4, Inches(9.5), Inches(0.1), Inches(3.5), Inches(0.276))
    
    # 统计表格数据
    stats_headers = ["九宫格", "指标", "已打", "备注"]
    for i, header in enumerate(stats_headers):
        cell = stats_table.table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = tencent_blue_medium
        cell.text_frame.paragraphs[0].font.color.rgb = white_color
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    
    # 使用动态计算的统计数据
    page3_stats_df = create_stats_table(df)
    stats_data = []
    for _, row in page3_stats_df.iterrows():
        stats_data.append([
            row['九宫格'],
            str(row['指标']),
            str(row['已打']),
            row['备注']
        ])
    
    for i, row_data in enumerate(stats_data, 1):
        for j, cell_data in enumerate(row_data):
            cell = stats_table.table.cell(i, j)
            cell.text = cell_data
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
            
            # 黄色高亮
            if j == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = tencent_yellow
    
    # 创建3x3九宫格
    grid_layout = [[7, 8, 9], [4, 5, 6], [1, 2, 3]]
    
    # 网格参数 - 为坐标轴留出空间
    start_left = Inches(1.2)  # 为Y轴留出空间
    start_top = Inches(1)
    cell_width = Inches(3.8)  # 稍微缩小以适应坐标轴
    cell_height = Inches(1.8)
    grid_width = cell_width * 3
    grid_height = cell_height * 3
    
    # 添加Y轴 (潜力)
    # Y轴线
    y_axis_line = slide.shapes.add_connector(
        1,  # 直线连接器
        Inches(0.8), start_top + grid_height,  # 起点：左下角
        Inches(0.8), start_top - Inches(0.3)   # 终点：向上延伸
    )
    y_axis_line.line.color.rgb = tencent_blue_standard
    y_axis_line.line.width = Pt(2)
    
    # Y轴箭头
    y_arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(0.7), start_top - Inches(0.4),
        Inches(0.2), Inches(0.2)
    )
    y_arrow.fill.solid()
    y_arrow.fill.fore_color.rgb = tencent_blue_standard
    y_arrow.line.color.rgb = tencent_blue_standard
    y_arrow.rotation = 90  # 旋转90度指向上方
    
    # Y轴标签
    y_label_box = slide.shapes.add_textbox(
        Inches(0.3), start_top + grid_height/2 - Inches(0.2),
        Inches(0.4), Inches(0.4)
    )
    y_label_frame = y_label_box.text_frame
    y_label_p = y_label_frame.paragraphs[0]
    y_label_p.text = "潜力"
    y_label_p.font.size = Pt(14)
    y_label_p.font.bold = True
    y_label_p.font.color.rgb = tencent_blue_standard
    y_label_p.font.name = "Microsoft YaHei"
    
    # 添加X轴 (绩效)
    # X轴线
    x_axis_line = slide.shapes.add_connector(
        1,  # 直线连接器
        Inches(0.8), start_top + grid_height,  # 起点：左下角
        start_left + grid_width + Inches(0.3), start_top + grid_height  # 终点：向右延伸
    )
    x_axis_line.line.color.rgb = tencent_blue_standard
    x_axis_line.line.width = Pt(2)
    
    # X轴箭头
    x_arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        start_left + grid_width + Inches(0.3), start_top + grid_height - Inches(0.1),
        Inches(0.2), Inches(0.2)
    )
    x_arrow.fill.solid()
    x_arrow.fill.fore_color.rgb = tencent_blue_standard
    x_arrow.line.color.rgb = tencent_blue_standard
    
    # X轴标签
    x_label_box = slide.shapes.add_textbox(
        start_left + grid_width/2 - Inches(0.2), start_top + grid_height + Inches(0.1),
        Inches(0.4), Inches(0.3)
    )
    x_label_frame = x_label_box.text_frame
    x_label_p = x_label_frame.paragraphs[0]
    x_label_p.text = "绩效"
    x_label_p.font.size = Pt(14)
    x_label_p.font.bold = True
    x_label_p.font.color.rgb = tencent_blue_standard
    x_label_p.font.name = "Microsoft YaHei"
    
    for row_idx, row in enumerate(grid_layout):
        for col_idx, rating in enumerate(row):
            # 计算位置
            left = start_left + col_idx * cell_width
            top = start_top + row_idx * cell_height
            
            # 添加带边框的形状
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, cell_width, cell_height
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(248, 249, 250)  # 浅灰背景
            rect.line.color.rgb = tencent_blue_medium
            rect.line.width = Pt(2)
            
            # 添加文本框
            textbox = slide.shapes.add_textbox(left, top, cell_width, cell_height)
            text_frame = textbox.text_frame
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            # 获取该档位的信息
            summary, count = generate_summary(df, rating)
            potential_level = get_potential_level(rating)
            
            # 标题段落
            p = text_frame.paragraphs[0]
            p.text = f"{rating} {potential_level}，低绩效-{count}人"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = tencent_blue_standard
            p.font.name = "Microsoft YaHei"
            
            # 内容段落 - 特殊处理档位5
            if rating == 5:
                # 档位5显示"详细名单请见下一页"
                p = text_frame.add_paragraph()
                p.text = "详细名单请见下一页"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = tencent_blue_medium
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER
            else:
                # 其他档位正常显示人员名单
                if summary and summary != "无符合员工":
                    content_lines = summary.split('\n')
                    for line in content_lines:
                        if line.strip():
                            p = text_frame.add_paragraph()
                            p.text = line.strip()
                            p.font.size = Pt(9)
                            p.font.name = "Microsoft YaHei"
                else:
                    p = text_frame.add_paragraph()
                    p.text = "无符合员工"
                    p.font.size = Pt(10)
                    p.font.italic = True
                    p.font.name = "Microsoft YaHei"
    
    # 添加绩效箭头
    arrow_box = slide.shapes.add_textbox(Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.4))
    arrow_frame = arrow_box.text_frame
    arrow_p = arrow_frame.paragraphs[0]
    arrow_p.text = "绩效"
    arrow_p.font.size = Pt(14)
    arrow_p.font.bold = True
    arrow_p.font.name = "Microsoft YaHei"
    arrow_p.alignment = PP_ALIGN.CENTER
    
    # 第四页：九宫格5档位详细名单
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 获取档位5的人员数据
    rating_5_df = df[df['档位'] == 5].copy() if '档位' in df.columns else pd.DataFrame()
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = f"附：九宫格5名单一览（{len(rating_5_df)}人）"
    title_p.font.size = Pt(24)
    title_p.font.bold = True
    title_p.font.color.rgb = tencent_blue_standard
    title_p.font.name = "Microsoft YaHei"
    
    if not rating_5_df.empty:
        # 按部门分组档位5的数据
        dept_groups = rating_5_df.groupby('部门')['姓名'].apply(list).to_dict()
        
        # 准备表格数据，保持部门分组结构
        all_data = []
        for dept, names in dept_groups.items():
            all_data.append({'dept': dept, 'names': names, 'count': len(names)})
        
        # 计算总人数并分成三个表格
        total_people = sum(item['count'] for item in all_data)
        target_per_table = total_people // 3
        
        # 智能分配部门到三个表格
        table1_data = []
        table2_data = []
        table3_data = []
        table1_count = 0
        table2_count = 0
        
        for item in all_data:
            if table1_count + item['count'] <= target_per_table or not table1_data:
                table1_data.append(item)
                table1_count += item['count']
            elif table2_count + item['count'] <= target_per_table or not table2_data:
                table2_data.append(item)
                table2_count += item['count']
            else:
                table3_data.append(item)
        
        # 定义表格创建函数
        def create_table(table_data, left_pos, table_name):
            if not table_data:
                return
                
            # 计算表格行数
            rows = sum(len(item['names']) for item in table_data) + 1
            
            table = slide.shapes.add_table(
                rows, 2,
                left_pos, Inches(1.2),
                Inches(4.2), Inches(5.5)
            )
            
            # 设置表头
            table.table.cell(0, 0).text = "部门"
            table.table.cell(0, 1).text = "员工姓名"
            
            # 表头样式
            for i in range(2):
                cell = table.table.cell(0, i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = tencent_blue_medium
                cell.text_frame.paragraphs[0].font.color.rgb = white_color
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 填充数据并合并部门单元格
            current_row = 1
            for dept_item in table_data:
                dept_name = dept_item['dept']
                names = dept_item['names']
                dept_start_row = current_row
                
                # 填充该部门的所有员工
                for name in names:
                    if current_row < len(table.table.rows):
                        table.table.cell(current_row, 0).text = dept_name if current_row == dept_start_row else ""
                        table.table.cell(current_row, 1).text = name
                        
                        # 设置样式
                        for j in range(2):
                            cell = table.table.cell(current_row, j)
                            cell.text_frame.paragraphs[0].font.size = Pt(8)
                            cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
                            if j == 0:
                                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    current_row += 1
                
                # 合并部门单元格
                if len(names) > 1 and dept_start_row < len(table.table.rows):
                    dept_end_row = min(current_row - 1, len(table.table.rows) - 1)
                    if dept_end_row > dept_start_row:
                        table.table.cell(dept_start_row, 0).merge(table.table.cell(dept_end_row, 0))
        
        # 创建三个并列表格
        create_table(table1_data, Inches(0.3), "左表格")    # 左表格
        create_table(table2_data, Inches(4.8), "中表格")    # 中表格  
        create_table(table3_data, Inches(9.3), "右表格")    # 右表格
    else:
        # 如果没有档位5的数据，显示提示信息
        no_data_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
        no_data_frame = no_data_box.text_frame
        no_data_p = no_data_frame.paragraphs[0]
        no_data_p.text = "暂无档位5的人员数据"
        no_data_p.font.size = Pt(18)
        no_data_p.font.color.rgb = tencent_blue_medium
        no_data_p.font.name = "Microsoft YaHei"
        no_data_p.alignment = PP_ALIGN.CENTER
    
    # 第五页：梯队重点关注
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 获取梯队重点关注数据
    cadre_not_tiered, expert_bottom5, expert_not_tiered, title_text = get_tier_focus_data(df)
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "附：梯队重点关注 | "
    title_p.font.size = Pt(20)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(220, 20, 60)  # 红色
    title_p.font.name = "Microsoft YaHei"
    
    # 添加统计说明文字
    title_p.text += title_text
    
    # 创建表格函数
    def create_focus_table(slide, data, title, left, top, width, height, headers):
        if data.empty:
            # 如果没有数据，创建空表格
            table = slide.shapes.add_table(6, len(headers), left, top, width, height)
        else:
            # 根据数据量决定表格行数，最少6行
            rows = max(6, len(data) + 1)
            table = slide.shapes.add_table(rows, len(headers), left, top, width, height)
        
        # 设置表头
        for i, header in enumerate(headers):
            cell = table.table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = tencent_blue_medium
            cell.text_frame.paragraphs[0].font.color.rgb = white_color
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 对数据按负责人和部门排序
        if not data.empty:
            data_sorted = data.sort_values(['负责人', '部门'], na_position='last')
        else:
            data_sorted = data
        
        # 填充数据
        for i, (_, row) in enumerate(data_sorted.iterrows(), 1):
            if i >= len(table.table.rows):
                break
                
            # 填充负责人信息（从Excel数据中提取）
            table.table.cell(i, 0).text = str(row.get('负责人', ''))
            table.table.cell(i, 1).text = str(row.get('部门', ''))
            table.table.cell(i, 2).text = str(row.get('姓名', ''))
            
            # 根据表格类型填充职级信息
            if '管理职级' in headers:
                table.table.cell(i, 3).text = str(row.get('管理职级', ''))
            else:
                # 使用优化的专业职级显示格式
                table.table.cell(i, 3).text = extract_professional_display(row.get('专业职级', ''))
            
            # 填充备注信息（如果Excel中有备注列）
            table.table.cell(i, 4).text = str(row.get('备注', ''))
            
            # 设置数据行样式
            for j in range(len(headers)):
                cell = table.table.cell(i, j)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
        
        return table
        
        return table
    
    # 左上表格：基干未盘入梯队
    headers_cadre = ["负责人", "部门", "员工姓名", "管理职级", "备注"]
    create_focus_table(
        slide, cadre_not_tiered, "基干未盘入梯队",
        Inches(0.5), Inches(1.0), Inches(6), Inches(2.8),
        headers_cadre
    )
    
    # 左下表格：专家末5%名单
    headers_expert = ["负责人", "部门", "员工姓名", "专业职级", "备注"]
    create_focus_table(
        slide, expert_bottom5, "专家末5%名单",
        Inches(0.5), Inches(4.0), Inches(6), Inches(2.8),
        headers_expert
    )
    
    # 右侧大表格：专家未盘入梯队人员名单
    right_table = create_focus_table(
        slide, expert_not_tiered, "专家未盘入梯队人员名单",
        Inches(7.0), Inches(1.0), Inches(6), Inches(5.8),
        headers_expert
    )
    
    # 设置右侧表格的行高和列宽
    if right_table:
        # 设置行高为0.4厘米 (约0.157英寸)
        for row in right_table.table.rows:
            row.height = Inches(0.157)
        
        # 设置列宽：负责人(3.2cm), 部门(3.8cm), 员工姓名(5cm), 专业职级(2.2cm), 备注(1.5cm)
        if len(right_table.table.columns) >= 5:
            right_table.table.columns[0].width = Inches(1.26)   # 负责人 3.2cm
            right_table.table.columns[1].width = Inches(1.50)   # 部门 3.8cm  
            right_table.table.columns[2].width = Inches(1.97)   # 员工姓名 5cm
            right_table.table.columns[3].width = Inches(0.87)   # 专业职级 2.2cm
            right_table.table.columns[4].width = Inches(0.59)   # 备注 1.5cm
        
        # 设置所有单元格的上下边距为0
        for row in right_table.table.rows:
            for cell in row.cells:
                cell.margin_top = 0
                cell.margin_bottom = 0
    
    # 添加表格标题
    # 左上表格标题
    title1_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(6), Inches(0.2))
    title1_frame = title1_box.text_frame
    title1_p = title1_frame.paragraphs[0]
    title1_p.text = "基干未盘入梯队"
    title1_p.font.size = Pt(12)
    title1_p.font.bold = True
    title1_p.font.color.rgb = tencent_blue_standard
    title1_p.font.name = "Microsoft YaHei"
    
    # 左下表格标题
    title2_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(6), Inches(0.2))
    title2_frame = title2_box.text_frame
    title2_p = title2_frame.paragraphs[0]
    title2_p.text = "专家末5%名单"
    title2_p.font.size = Pt(12)
    title2_p.font.bold = True
    title2_p.font.color.rgb = tencent_blue_standard
    title2_p.font.name = "Microsoft YaHei"
    
    # 右侧表格标题
    title3_box = slide.shapes.add_textbox(Inches(7.0), Inches(0.8), Inches(6), Inches(0.2))
    title3_frame = title3_box.text_frame
    title3_p = title3_frame.paragraphs[0]
    title3_p.text = "专家未盘入梯队人员名单"
    title3_p.font.size = Pt(12)
    title3_p.font.bold = True
    title3_p.font.color.rgb = tencent_blue_standard
    title3_p.font.name = "Microsoft YaHei"
    

    
    return prs

def main():
    """
    主应用函数
    """
    # 页面标题
    st.markdown('<div class="main-title">📊 S1人才盘点系统</div>', 
                unsafe_allow_html=True)
    
    # 侧边栏 - 文件上传
    st.sidebar.header("📁 数据上传")
    
    # 提供示例数据选项
    use_sample = st.sidebar.checkbox("使用示例数据进行测试")
    
    df = None
    
    if use_sample:
        df = create_sample_data()
        st.sidebar.success("已加载示例数据")
    else:
        uploaded_file = st.sidebar.file_uploader(
            "上传Excel文件", 
            type=['xlsx', 'xls'],
            help="请上传包含'九宫格'、'员工姓名'、'部门'、'梯队'列的Excel文件"
        )
        
        if uploaded_file is not None:
            sheet_name = st.sidebar.text_input("工作表名称", value="Sheet1")
            df = load_and_validate_data(uploaded_file, sheet_name)
    
    if df is not None and not df.empty:
        # 数据概览
        total_count = len(df)
        
        # 1 数据总览区域
        st.markdown('<div class="main-title">1 数据总览</div>', unsafe_allow_html=True)
        
        # 创建左右两列布局
        col1, col2 = st.columns(2)
        
        with col1:
            # 梯队统计表格 (左侧)
            st.markdown("**梯队总览**")
            

            
            tier_summary = generate_tier_summary(df)
            st.markdown(f"<div style='font-size: 12px; color: #666; margin-bottom: 10px;'>{tier_summary}</div>", 
                       unsafe_allow_html=True)
            
            tier_stats_df = create_tier_stats_table(df)
            if not tier_stats_df.empty:
                st.markdown('<div class="stats-table">', unsafe_allow_html=True)
                st.dataframe(tier_stats_df, width='stretch', hide_index=True)
                st.markdown('</div>', unsafe_allow_html=True)
        
        with col2:
            # 九宫格统计表格 (右侧)
            st.markdown("**基干九宫格**")
            grid_summary = generate_grid_summary(df)
            st.markdown(f"<div style='font-size: 12px; color: #666; margin-bottom: 10px;'>{grid_summary}</div>", 
                       unsafe_allow_html=True)
            
            stats_df = create_stats_table(df)
            if not stats_df.empty:
                st.markdown('<div class="stats-table">', unsafe_allow_html=True)
                st.dataframe(stats_df, width='stretch', hide_index=True)
                st.markdown('</div>', unsafe_allow_html=True)
        
        # 九宫格布局
        # 计算有九宫格数据的人数
        grid_count = len(df[df['档位'].notna()])
        st.markdown(f'<div class="main-title">2 基干九宫格一览：总计{grid_count}人</div>', 
                   unsafe_allow_html=True)
        st.markdown("<br>", unsafe_allow_html=True)
        
        # 定义九宫格排列 (7,8,9 / 4,5,6 / 1,2,3)
        grid_layout = [
            [7, 8, 9],
            [4, 5, 6], 
            [1, 2, 3]
        ]
        
        # 创建3行布局
        for row_idx, row in enumerate(grid_layout):
            cols = st.columns(3)
            
            for i, rating in enumerate(row):
                with cols[i]:
                    # 获取该档位的汇总信息
                    summary, count = generate_summary(df, rating)
                    potential_level = get_potential_level(rating)
                    
                    # 格子标题 - 显示实际人数
                    title = f"{rating} {potential_level}，低绩效-{count}人"
                    
                    # 使用HTML容器显示内容
                    if summary and summary != "无符合员工":
                        # 确保HTML转义和格式正确
                        content_lines = []
                        for line in summary.split('\n'):
                            if line.strip():
                                # HTML转义特殊字符
                                escaped_line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                                content_lines.append(escaped_line)
                        content = "<br>".join(content_lines)
                    else:
                        content = "无符合员工"
                    
                    # 使用完整的HTML容器
                    cell_html = f"""
                    <div class="grid-cell">
                        <div class="cell-title">{title}</div>
                        <div class="cell-content">{content}</div>
                    </div>
                    """
                    st.markdown(cell_html, unsafe_allow_html=True)
            
            # 在每行后添加间距
            st.markdown("<br>", unsafe_allow_html=True)
        
        # 第三部分：梯队重点关注
        st.markdown('<div class="main-title">3 梯队重点关注</div>', unsafe_allow_html=True)
        st.markdown("<br>", unsafe_allow_html=True)
        
        # 获取梯队重点关注数据
        cadre_not_tiered, expert_bottom5, expert_not_tiered, title_text = get_tier_focus_data(df)
        
        # 显示统计摘要
        st.markdown(f'<div style="color: #DC143C; font-size: 16px; font-weight: bold; margin-bottom: 20px;">附：梯队重点关注 | {title_text}</div>', 
                   unsafe_allow_html=True)
        
        # 创建三栏布局
        col1, col2 = st.columns([1, 1])
        
        with col1:
            # 左上：基干未盘入梯队
            st.markdown("**基干未盘入梯队**")
            if not cadre_not_tiered.empty:
                # 对数据排序并优化显示
                cadre_display = cadre_not_tiered.sort_values(['负责人', '部门'], na_position='last')
                # 选择存在的列
                display_cols = ['负责人', '部门', '姓名', '管理职级']
                if '备注' in cadre_display.columns:
                    display_cols.append('备注')
                cadre_display = cadre_display[display_cols].copy()
                st.dataframe(cadre_display, width='stretch', hide_index=True)
            else:
                st.info("暂无基干未盘入梯队人员")
            
            st.markdown("<br>", unsafe_allow_html=True)
            
            # 左下：专家末5%名单
            st.markdown("**专家末5%名单**")
            if not expert_bottom5.empty:
                # 对数据排序并优化专业职级显示
                expert5_display = expert_bottom5.sort_values(['负责人', '部门'], na_position='last').copy()
                expert5_display['专业职级显示'] = expert5_display['专业职级'].apply(extract_professional_display)
                # 选择存在的列
                display_cols = ['负责人', '部门', '姓名', '专业职级显示']
                if '备注' in expert5_display.columns:
                    display_cols.append('备注')
                expert5_display = expert5_display[display_cols].copy()
                expert5_display = expert5_display.rename(columns={'专业职级显示': '专业职级'})
                st.dataframe(expert5_display, width='stretch', hide_index=True)
            else:
                st.info("暂无专家末5%人员")
        
        with col2:
            # 右侧：专家未盘入梯队人员名单
            st.markdown("**专家未盘入梯队人员名单**")
            if not expert_not_tiered.empty:
                # 对数据排序并优化专业职级显示
                expert_display = expert_not_tiered.sort_values(['负责人', '部门'], na_position='last').copy()
                expert_display['专业职级显示'] = expert_display['专业职级'].apply(extract_professional_display)
                # 选择存在的列
                display_cols = ['负责人', '部门', '姓名', '专业职级显示']
                if '备注' in expert_display.columns:
                    display_cols.append('备注')
                expert_display = expert_display[display_cols].copy()
                expert_display = expert_display.rename(columns={'专业职级显示': '专业职级'})
                st.dataframe(expert_display, width='stretch', hide_index=True, height=400)
            else:
                st.info("暂无专家未盘入梯队人员")
        
        st.markdown("<br><br>", unsafe_allow_html=True)
        
        # 数据详情
        with st.expander("📋 查看原始数据"):
            st.dataframe(df, width='stretch')
            
        # 下载处理后的数据
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("📥 下载Excel报告"):
                # 创建汇总报告
                report_data = []
                for rating in range(1, 10):
                    summary, count = generate_summary(df, rating)
                    potential_level = get_potential_level(rating)
                    report_data.append({
                        '档位': rating,
                        '潜力级别': potential_level,
                        '人数': count,
                        '员工详情': summary.replace('\n', '; ') if summary != "无符合员工" else "无"
                    })
                
                report_df = pd.DataFrame(report_data)
                
                # 转换为Excel
                output = BytesIO()
                with pd.ExcelWriter(output, engine='openpyxl') as writer:
                    report_df.to_excel(writer, sheet_name='九宫格汇总', index=False)
                    df.to_excel(writer, sheet_name='原始数据', index=False)
                
                st.download_button(
                    label="下载Excel文件",
                    data=output.getvalue(),
                    file_name="九宫格潜力汇总报告.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
        
        with col2:
            if st.button("📊 下载PPT报告"):
                try:
                    # 创建PPT报告
                    ppt = create_ppt_report(df)
                    
                    # 保存到BytesIO
                    ppt_output = BytesIO()
                    ppt.save(ppt_output)
                    ppt_output.seek(0)
                    
                    st.download_button(
                        label="下载PPT文件",
                        data=ppt_output.getvalue(),
                        file_name="九宫格潜力展示报告.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    
                    st.success("✅ PPT报告生成成功！包含以下内容：")
                    st.info("""
                    📋 PPT内容包括：
                    • 第1页：标题页和总体统计
                    • 第2页：潜力分布统计概览  
                    • 第3页：九宫格详细分布图
                    • 第4页：完整员工数据表
                    
                    💡 所有文字内容均可在PowerPoint中编辑
                    """)
                    
                except Exception as e:
                    st.error(f"生成PPT时出错：{str(e)}")
    
    else:
        # 显示使用说明
        st.info("👆 请在左侧上传Excel文件或选择使用示例数据")
        
        st.markdown("""
        ### 📋 使用说明
        
        1. **数据格式要求**：
           - Excel文件需包含四列：`九宫格`、`员工姓名`、`部门`、`梯队`
           - 九宫格列应为1-9的数字
           - 员工姓名和部门列为文本格式
           - 梯队列应为：前5%、前15%、前40%、末5%
        
        2. **九宫格说明**：
           - 7,8,9 → 高潜力
           - 4,5,6 → 中潜力  
           - 1,2,3 → 低潜力
        
        3. **梯队说明**：
           - 前5% → 顶尖人才
           - 前15% → 核心人才
           - 前40% → 骨干人才
           - 末5% → 重点关注
        
        4. **功能特点**：
           - 双维度人才分析（九宫格+梯队）
           - 自动按部门汇总员工信息
           - 实时统计各档位和梯队人数
           - 支持S1人才盘点PPT报告生成
        """)

if __name__ == "__main__":
    main()